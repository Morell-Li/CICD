#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build Morell 述职总结 — 科技感深色版 PPT.

Run:
    .venv-ppt/bin/python scripts/build_morell_deck.py
"""
from __future__ import annotations

import os
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


OUT_PATH = os.path.expanduser("~/Desktop/Morell述职总结_美化版.pptx")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# Design tokens
# ============================================================
class C:
    BG_DEEP    = RGBColor(0x0B, 0x14, 0x26)
    BG_DEEP_2  = RGBColor(0x0E, 0x1B, 0x2E)
    BG_PANEL   = RGBColor(0x15, 0x22, 0x38)
    BG_PANEL_2 = RGBColor(0x1B, 0x2C, 0x47)
    BG_PANEL_3 = RGBColor(0x0F, 0x1C, 0x33)
    ACCENT     = RGBColor(0x00, 0xE5, 0xFF)
    ACCENT_2   = RGBColor(0x4D, 0x9F, 0xFF)
    ACCENT_3   = RGBColor(0x7C, 0xFF, 0xE0)
    GLOW       = RGBColor(0x0E, 0x25, 0x45)
    TEXT_PRI   = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_SEC   = RGBColor(0xB8, 0xC5, 0xD1)
    TEXT_DIM   = RGBColor(0x6B, 0x7A, 0x8F)
    LINE       = RGBColor(0x2A, 0x3D, 0x5C)
    DIGIT_GHOST = RGBColor(0x18, 0x2C, 0x55)


FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"


# ============================================================
# Low-level helpers
# ============================================================
def _set_east_asian_font(run, font_name: str):
    """Force CJK glyphs to the given typeface."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:latin", "a:cs"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font_name)


def _remove_shadow(shape):
    sp = shape._element
    for el in sp.findall(".//" + qn("a:effectLst")):
        el.getparent().remove(el)


def _set_shape_fill_alpha(shape, alpha_pct: float):
    """Add alpha to a solidFill via XML. alpha_pct 0-100."""
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return
    color_el = None
    for tag in ("a:srgbClr", "a:schemeClr"):
        color_el = solid.find(qn(tag))
        if color_el is not None:
            break
    if color_el is None:
        return
    for a in color_el.findall(qn("a:alpha")):
        color_el.remove(a)
    alpha = etree.SubElement(color_el, qn("a:alpha"))
    alpha.set("val", str(int(alpha_pct * 1000)))


def set_bg(slide, color=C.BG_DEEP):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, *,
             fill: Optional[RGBColor] = None,
             line: Optional[RGBColor] = None,
             line_w=None,
             no_line=False,
             kind=MSO_SHAPE.RECTANGLE,
             rounded_adj: Optional[float] = None,
             alpha: Optional[float] = None):
    sh = slide.shapes.add_shape(kind, left, top, w, h)
    if rounded_adj is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = rounded_adj
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if alpha is not None:
            _set_shape_fill_alpha(sh, alpha)
    if no_line:
        sh.line.fill.background()
    elif line is not None:
        sh.line.color.rgb = line
        if line_w is not None:
            sh.line.width = line_w
    sh.shadow.inherit = False
    _remove_shadow(sh)
    return sh


def add_text(slide, left, top, w, h, text, *,
             size=18, color=C.TEXT_PRI, bold=False, italic=False,
             align=PP_ALIGN.LEFT, vert=MSO_ANCHOR.TOP,
             font=FONT_CN, line_space: Optional[float] = None,
             letter_space: Optional[int] = None):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = vert
    p = tf.paragraphs[0]
    p.alignment = align
    if line_space is not None:
        p.line_spacing = line_space
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    _set_east_asian_font(r, font)
    if letter_space is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(letter_space))
    return tb


def add_runs(slide, left, top, w, h, paragraphs, *,
             vert=MSO_ANCHOR.TOP, default_align=PP_ALIGN.LEFT,
             line_space: float = 1.35):
    """
    paragraphs: list[ dict(align?, space_before?, line_space?, runs=[ dict(text,size,color,bold,italic,font) ]) ]
    """
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = vert
    for i, par in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = par.get("align", default_align)
        p.line_spacing = par.get("line_space", line_space)
        sb = par.get("space_before")
        if sb is not None:
            p.space_before = Pt(sb)
        sa = par.get("space_after")
        if sa is not None:
            p.space_after = Pt(sa)
        for r in par["runs"]:
            run = p.add_run()
            run.text = r["text"]
            f = run.font
            f.size = Pt(r.get("size", 18))
            f.bold = r.get("bold", False)
            f.italic = r.get("italic", False)
            f.color.rgb = r.get("color", C.TEXT_PRI)
            font_name = r.get("font", FONT_CN)
            f.name = font_name
            _set_east_asian_font(run, font_name)
    return tb


# ============================================================
# Page chrome
# ============================================================
def add_top_accent(slide, *, color=C.ACCENT, height_pt=2):
    add_rect(slide, 0, 0, SLIDE_W, Pt(height_pt), fill=color, no_line=True)


def add_corner_brand(slide):
    add_rect(slide, Inches(0.55), Inches(0.32), Inches(0.18), Pt(2),
             fill=C.ACCENT, no_line=True)
    add_text(slide, Inches(0.8), Inches(0.22), Inches(6), Inches(0.4),
             "MORELL  ·  述职总结  ·  2026/05/12",
             size=10, color=C.TEXT_DIM, font=FONT_EN, letter_space=200)


def add_page_number(slide, idx, total=15, *, color=C.TEXT_DIM):
    add_text(slide, Inches(12.0), Inches(7.05), Inches(1.1), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=10, color=color,
             align=PP_ALIGN.RIGHT, font=FONT_EN, letter_space=200)


def add_footer_line(slide):
    add_rect(slide, Inches(0.55), Inches(7.08), Inches(0.4), Pt(1.5),
             fill=C.ACCENT, no_line=True)


def add_section_eyebrow(slide, eyebrow_en, title_cn):
    add_text(slide, Inches(0.55), Inches(0.85), Inches(8), Inches(0.3),
             eyebrow_en, size=11, color=C.ACCENT, bold=True,
             font=FONT_EN, letter_space=400)
    add_text(slide, Inches(0.55), Inches(1.15), Inches(10), Inches(0.7),
             title_cn, size=30, color=C.TEXT_PRI, bold=True)
    add_rect(slide, Inches(0.55), Inches(1.95), Inches(0.55), Pt(2.5),
             fill=C.ACCENT, no_line=True)


def add_card(slide, left, top, w, h, *,
             fill=C.BG_PANEL, line=C.LINE, line_w=Pt(0.75),
             rounded_adj=0.05):
    return add_rect(slide, left, top, w, h, fill=fill, line=line,
                    line_w=line_w, kind=MSO_SHAPE.ROUNDED_RECTANGLE,
                    rounded_adj=rounded_adj)


def add_card_accent_bar(slide, left, top, w=Inches(0.06), h=Inches(0.6),
                        color=C.ACCENT):
    add_rect(slide, left, top, w, h, fill=color, no_line=True)


def add_pill(slide, left, top, text, *, fill=C.BG_PANEL_2, color=C.ACCENT):
    char_w = 0.155
    pad = 0.22
    w = Inches(char_w * len(text) + pad * 2)
    h = Inches(0.36)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.adjustments[0] = 0.5
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = color
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    _remove_shadow(sh)
    _set_shape_fill_alpha(sh, 60)
    tf = sh.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_CN
    r.font.size = Pt(11)
    r.font.color.rgb = color
    r.font.bold = False
    _set_east_asian_font(r, FONT_CN)
    return sh, w


def add_grid_decor(slide, origin_left, origin_top, *, cols=10, rows=5,
                   cell=Inches(0.2), color=C.ACCENT_2, alpha=22):
    """Decorative dot grid bottom-right."""
    for r in range(rows):
        for c in range(cols):
            dot = add_rect(
                slide,
                origin_left + c * cell,
                origin_top + r * cell,
                Inches(0.04),
                Inches(0.04),
                fill=color,
                no_line=True,
            )
            _set_shape_fill_alpha(dot, alpha)


def add_glow_circle(slide, cx_in, cy_in, r_in, *, color=C.ACCENT, alpha=12):
    sh = add_rect(slide,
                  Inches(cx_in - r_in), Inches(cy_in - r_in),
                  Inches(r_in * 2), Inches(r_in * 2),
                  fill=color, no_line=True, kind=MSO_SHAPE.OVAL)
    _set_shape_fill_alpha(sh, alpha)
    return sh


def add_ghost_number(slide, left, top, w, h, num, *, size=240,
                     color=C.DIGIT_GHOST):
    return add_text(slide, left, top, w, h, num, size=size, color=color,
                    bold=True, font=FONT_EN, align=PP_ALIGN.LEFT,
                    vert=MSO_ANCHOR.MIDDLE)


def page_chrome(slide, idx):
    add_top_accent(slide)
    add_corner_brand(slide)
    add_page_number(slide, idx)
    add_footer_line(slide)


# ============================================================
# Slides
# ============================================================
def slide_1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # decorative glows
    add_glow_circle(s, 12.5, 6.5, 2.2, color=C.ACCENT, alpha=8)
    add_glow_circle(s, 0.8, 0.8, 1.5, color=C.ACCENT_2, alpha=10)
    add_grid_decor(s, Inches(10.2), Inches(5.6), cols=14, rows=6,
                   cell=Inches(0.22), color=C.ACCENT, alpha=20)
    # top bar
    add_top_accent(s, height_pt=3)
    # date eyebrow
    add_text(s, Inches(0.9), Inches(0.7), Inches(6), Inches(0.4),
             "2026 / 05 / 12   ·   ANNUAL REVIEW",
             size=12, color=C.ACCENT, bold=True, font=FONT_EN,
             letter_space=400)
    # short accent bar above title
    add_rect(s, Inches(0.9), Inches(1.5), Inches(0.6), Pt(3),
             fill=C.ACCENT, no_line=True)
    # main title
    add_runs(s, Inches(0.9), Inches(1.85), Inches(11), Inches(2.2), [
        {"runs": [{"text": "Morell ", "size": 64, "bold": True,
                   "color": C.TEXT_PRI, "font": FONT_EN},
                  {"text": "述职总结", "size": 60, "bold": True,
                   "color": C.TEXT_PRI}],
         "line_space": 1.0},
    ])
    # subtitle accent
    add_text(s, Inches(0.9), Inches(3.95), Inches(11), Inches(0.6),
             "A Year of Building · Testing · Shipping",
             size=18, color=C.ACCENT_2, italic=True, font=FONT_EN,
             letter_space=150)
    # divider
    add_rect(s, Inches(0.9), Inches(4.85), Inches(2.8), Pt(1.5),
             fill=C.LINE, no_line=True)
    # presenter
    add_runs(s, Inches(0.9), Inches(5.1), Inches(8), Inches(0.6), [
        {"runs": [
            {"text": "汇报人  ", "size": 14, "color": C.TEXT_DIM,
             "font": FONT_EN, "bold": True},
            {"text": "Morell", "size": 20, "color": C.TEXT_PRI,
             "bold": True, "font": FONT_EN},
        ]},
    ])
    # bottom-right brand strip
    add_text(s, Inches(0.9), Inches(6.95), Inches(6), Inches(0.3),
             "ZOOMEX  ·  QA  ·  BASE TRADING GROUP",
             size=10, color=C.TEXT_DIM, font=FONT_EN, letter_space=400)


def slide_2_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 2)
    # left
    add_text(s, Inches(0.6), Inches(1.4), Inches(5), Inches(0.4),
             "CONTENTS", size=14, color=C.ACCENT, bold=True,
             font=FONT_EN, letter_space=400)
    add_text(s, Inches(0.6), Inches(1.85), Inches(5), Inches(1.2),
             "目录", size=72, color=C.TEXT_PRI, bold=True)
    add_rect(s, Inches(0.6), Inches(3.4), Inches(0.6), Pt(2.5),
             fill=C.ACCENT, no_line=True)
    add_text(s, Inches(0.6), Inches(3.7), Inches(5), Inches(0.6),
             "Annual Self-Review · 2026",
             size=12, color=C.TEXT_DIM, font=FONT_EN, italic=True,
             letter_space=150)

    # right cards (3 vertical)
    sections = [
        ("01", "公司产品 & 文化理念",
         "Company products and culture"),
        ("02", "部门业务 & OKR",
         "Department business and OKR"),
        ("03", "个人总结",
         "Personal review"),
    ]
    card_left = Inches(6.4)
    card_w = Inches(6.3)
    card_h = Inches(1.35)
    gap = Inches(0.22)
    start_top = Inches(1.45)
    for i, (num, cn, en) in enumerate(sections):
        top = start_top + i * (card_h + gap)
        add_card(s, card_left, top, card_w, card_h,
                 fill=C.BG_PANEL, line=C.LINE)
        # accent vertical bar
        add_rect(s, card_left + Inches(0.18), top + Inches(0.25),
                 Pt(2.5), card_h - Inches(0.5), fill=C.ACCENT, no_line=True)
        # big number
        add_text(s, card_left + Inches(0.45), top + Inches(0.18),
                 Inches(1.4), Inches(1.0),
                 num, size=44, color=C.ACCENT, bold=True, font=FONT_EN,
                 vert=MSO_ANCHOR.MIDDLE)
        # title
        add_text(s, card_left + Inches(1.95), top + Inches(0.28),
                 card_w - Inches(2.2), Inches(0.6),
                 cn, size=22, color=C.TEXT_PRI, bold=True)
        add_text(s, card_left + Inches(1.95), top + Inches(0.78),
                 card_w - Inches(2.2), Inches(0.4),
                 en, size=11, color=C.TEXT_DIM, font=FONT_EN, italic=True,
                 letter_space=150)


def section_divider(prs, idx, num, cn_title, en_subtitle, page_idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, page_idx)
    # ghost giant number
    add_ghost_number(s, Inches(0.4), Inches(1.0), Inches(7), Inches(6),
                     num, size=420, color=C.DIGIT_GHOST)
    # accent vertical line on right side
    add_rect(s, Inches(7.5), Inches(2.4), Pt(2), Inches(2.6),
             fill=C.ACCENT, no_line=True)
    # eyebrow EN
    add_text(s, Inches(7.85), Inches(2.4), Inches(5), Inches(0.4),
             f"CHAPTER · {num}",
             size=12, color=C.ACCENT, bold=True, font=FONT_EN,
             letter_space=400)
    # CN title
    add_text(s, Inches(7.85), Inches(2.85), Inches(5.5), Inches(1.1),
             cn_title, size=42, color=C.TEXT_PRI, bold=True)
    # EN subtitle
    add_text(s, Inches(7.85), Inches(4.15), Inches(5.5), Inches(0.7),
             en_subtitle, size=14, color=C.TEXT_SEC, italic=True,
             font=FONT_EN, letter_space=150)
    # decorative dots
    add_grid_decor(s, Inches(10.5), Inches(5.6), cols=10, rows=4,
                   cell=Inches(0.2), color=C.ACCENT, alpha=18)
    add_glow_circle(s, 12.4, 1.6, 0.9, color=C.ACCENT_2, alpha=14)


def slide_4_culture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 4)
    add_section_eyebrow(s, "01 · COMPANY CULTURE", "公司文化")

    cards = [
        ("01", "Always day one", "保持创业的第一天",
         "把每一天都当作公司的第一天，保持饥饿与谦逊，持续学习、敢于试错。",
         C.ACCENT),
        ("02", "Respect every cent", "敬畏每一分钱",
         "对成本与资源高度敬畏，让每一分钱花在用户与产品价值最大化的地方。",
         C.ACCENT_2),
    ]
    card_w = Inches(5.85)
    card_h = Inches(3.95)
    left0 = Inches(0.6)
    gap = Inches(0.3)
    top = Inches(2.55)
    for i, (num, en, cn, desc, accent) in enumerate(cards):
        left = left0 + i * (card_w + gap)
        add_card(s, left, top, card_w, card_h)
        add_card_accent_bar(s, left + Inches(0.45), top + Inches(0.5),
                            w=Pt(3), h=Inches(0.7), color=accent)
        add_text(s, left + Inches(0.45), top + Inches(1.45),
                 Inches(3), Inches(0.4),
                 num, size=14, color=accent, bold=True, font=FONT_EN,
                 letter_space=300)
        add_text(s, left + Inches(0.45), top + Inches(1.75),
                 card_w - Inches(0.9), Inches(0.6),
                 en, size=24, color=C.TEXT_PRI, bold=True, font=FONT_EN)
        add_text(s, left + Inches(0.45), top + Inches(2.3),
                 card_w - Inches(0.9), Inches(0.6),
                 cn, size=20, color=accent, bold=True)
        add_text(s, left + Inches(0.45), top + Inches(2.95),
                 card_w - Inches(0.9), Inches(1.0),
                 desc, size=13, color=C.TEXT_SEC, line_space=1.6)


def _product_intro_slide(prs, page_idx, num, en_title, cn_title, url,
                         body_text, placeholder_label):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, page_idx)
    add_section_eyebrow(s, "01 · PRODUCTS", "公司产品介绍")

    # left text column
    left_l = Inches(0.6); left_w = Inches(6.6); top0 = Inches(2.55)
    add_text(s, left_l, top0, Inches(2), Inches(0.5),
             num, size=14, color=C.ACCENT, bold=True, font=FONT_EN,
             letter_space=400)
    add_text(s, left_l, top0 + Inches(0.45), left_w, Inches(0.8),
             en_title, size=30, color=C.TEXT_PRI, bold=True, font=FONT_EN)
    add_text(s, left_l, top0 + Inches(1.05), left_w, Inches(0.5),
             cn_title, size=18, color=C.ACCENT_2, bold=True)
    if url:
        add_text(s, left_l, top0 + Inches(1.6), left_w, Inches(0.4),
                 url, size=14, color=C.ACCENT, font=FONT_EN,
                 letter_space=100)
    add_rect(s, left_l, top0 + Inches(2.0), Inches(0.5), Pt(2),
             fill=C.LINE, no_line=True)
    add_text(s, left_l, top0 + Inches(2.2), left_w, Inches(2.5),
             body_text, size=15, color=C.TEXT_SEC, line_space=1.7)

    # right image placeholder
    p_left = Inches(7.65); p_top = Inches(2.45)
    p_w = Inches(5.1); p_h = Inches(4.4)
    add_card(s, p_left, p_top, p_w, p_h,
             fill=C.BG_PANEL_3, line=C.ACCENT, line_w=Pt(0.75),
             rounded_adj=0.04)
    # corner brackets to imply screenshot frame
    brk = Inches(0.35); bw = Pt(2.2)
    for (cx, cy, dx, dy) in [
        (p_left + Inches(0.15), p_top + Inches(0.15), 1, 1),
        (p_left + p_w - Inches(0.15) - brk, p_top + Inches(0.15), -1, 1),
        (p_left + Inches(0.15), p_top + p_h - Inches(0.15), 1, -1),
        (p_left + p_w - Inches(0.15) - brk, p_top + p_h - Inches(0.15), -1, -1),
    ]:
        add_rect(s, cx, cy, brk, bw, fill=C.ACCENT, no_line=True)
        add_rect(s, cx if dx > 0 else cx + brk - bw,
                 cy, bw, brk, fill=C.ACCENT, no_line=True)
    # center label
    add_text(s, p_left, p_top + p_h / 2 - Inches(0.5), p_w, Inches(0.4),
             "PRODUCT SCREENSHOT", size=11, color=C.ACCENT, bold=True,
             font=FONT_EN, align=PP_ALIGN.CENTER, letter_space=400)
    add_text(s, p_left, p_top + p_h / 2 - Inches(0.05), p_w, Inches(0.45),
             placeholder_label, size=14, color=C.TEXT_SEC,
             align=PP_ALIGN.CENTER)
    add_text(s, p_left, p_top + p_h / 2 + Inches(0.5), p_w, Inches(0.4),
             "在 Keynote / PowerPoint 中右键此框 → 替换为图片",
             size=10, color=C.TEXT_DIM, align=PP_ALIGN.CENTER)


def slide_5_product_web(prs):
    body = ("DEX & CEX 一体化交易平台。包含法币出入金、现货、合约、"
            "反向合约、闪兑、理财，KYC 认证以及 OpenAPI 接口等模块，"
            "面向 C 端用户与机构/做市商。")
    _product_intro_slide(prs, 5, "01", "WEB Zoomex", "网页端交易平台",
                         "www.zoomex.com", body, "Web 端首页 / 交易页")


def slide_6_product_app(prs):
    body = ("iOS & Android 双端原生 APP。包含出入金、现货、合约、反向合约、"
            "闪兑、理财，KYC 认证，资产账户等核心模块，移动场景下的全功能"
            "交易闭环。")
    _product_intro_slide(prs, 6, "02", "APP iOS & Android", "移动端 APP",
                         "", body, "APP 端关键页面")


def slide_8_dept(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 8)
    add_section_eyebrow(s, "02 · DEPARTMENT", "基础交易组 业务")

    # main card
    card_l = Inches(0.6); card_t = Inches(2.55)
    card_w = Inches(12.1); card_h = Inches(4.4)
    add_card(s, card_l, card_t, card_w, card_h)
    add_card_accent_bar(s, card_l + Inches(0.5), card_t + Inches(0.55),
                        w=Pt(3), h=Inches(0.7), color=C.ACCENT)
    add_text(s, card_l + Inches(0.5), card_t + Inches(0.45),
             Inches(8), Inches(0.5),
             "BUSINESS SCOPE", size=11, color=C.ACCENT, bold=True,
             font=FONT_EN, letter_space=400)
    add_text(s, card_l + Inches(0.5), card_t + Inches(0.85),
             Inches(11), Inches(0.8),
             "负责业务资产 · 财务 · 出入金 · 交易系统", size=26,
             color=C.TEXT_PRI, bold=True)
    add_text(s, card_l + Inches(0.5), card_t + Inches(1.7),
             Inches(11), Inches(1.0),
             "覆盖核心交易撮合（现货 / 合约 / 反向合约）与行情、KYC、"
             "Open API；OpenAPI 主要面向做市商和机构客户。",
             size=15, color=C.TEXT_SEC, line_space=1.7)

    # pills
    pills = ["业务资产", "财务", "出入金", "现货", "合约", "反向合约",
             "行情", "KYC", "OpenAPI", "做市商", "机构客户"]
    left = card_l + Inches(0.5)
    top = card_t + Inches(3.05)
    cur_left = left
    row_top = top
    max_right = card_l + card_w - Inches(0.5)
    for t in pills:
        sh, w = add_pill(s, cur_left, row_top, t)
        cur_left = cur_left + w + Inches(0.15)
        if cur_left + Inches(1.5) > max_right:
            cur_left = left
            row_top = row_top + Inches(0.5)


def slide_9_okr(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 9)
    add_section_eyebrow(s, "02 · OBJECTIVES", "OKR 目标拆解")

    okrs = [
        ("01", "网页版 APP (PWA)",
         "H5 站启动项目，完成合约、资产出入金、现货模块的开发上线。",
         "PWA · 合约 · 出入金 · 现货"),
        ("02", "交易 AI MCP / Agent",
         "开发 AI 交易工具，支持 AI 交易，提供「第四个交易端」，跟进主流交易所节奏。",
         "AI · MCP · Agent · 第四端"),
        ("03", "出入金流程优化",
         "完成所有版本的流程优化并全部上线主网。",
         "V1.0 · V2.0 · V3.0"),
        ("04", "测试自动化平台",
         "测试自动化平台应用于交易核心流程，建立可重复回归与持续运行机制。",
         "自动化 · 核心流程 · 持续回归"),
    ]
    grid_l = Inches(0.6); grid_t = Inches(2.55)
    card_w = Inches(5.95); card_h = Inches(2.1)
    gap_x = Inches(0.2); gap_y = Inches(0.2)
    for i, (num, t, desc, tags) in enumerate(okrs):
        row, col = divmod(i, 2)
        left = grid_l + col * (card_w + gap_x)
        top = grid_t + row * (card_h + gap_y)
        add_card(s, left, top, card_w, card_h)
        add_card_accent_bar(s, left + Inches(0.3), top + Inches(0.3),
                            w=Pt(2.5), h=Inches(0.6), color=C.ACCENT)
        add_text(s, left + Inches(0.5), top + Inches(0.22),
                 Inches(1.5), Inches(0.5),
                 num, size=22, color=C.ACCENT, bold=True, font=FONT_EN)
        add_text(s, left + Inches(1.4), top + Inches(0.25),
                 card_w - Inches(1.6), Inches(0.5),
                 t, size=18, color=C.TEXT_PRI, bold=True)
        add_text(s, left + Inches(0.5), top + Inches(0.95),
                 card_w - Inches(0.9), Inches(0.85),
                 desc, size=13, color=C.TEXT_SEC, line_space=1.6)
        add_text(s, left + Inches(0.5), top + Inches(1.65),
                 card_w - Inches(0.9), Inches(0.35),
                 tags, size=10, color=C.ACCENT_2, font=FONT_EN,
                 letter_space=200)


def slide_11_responsibilities(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 11)
    add_section_eyebrow(s, "03 · RESPONSIBILITIES", "岗位职责")

    items = [
        ("FUNCTIONAL", "功能测试",
         "负责基础交易组相关业务需求的需求拆解、用例编写、测试执行，"
         "并跟进问题闭环，保障交付质量。"),
        ("PERFORMANCE", "性能测试",
         "如有需求需要性能压测，负责对应服务的性能指标压测，"
         "压测方案梳理与执行，输出性能报告。"),
        ("AUTOMATION", "自动化测试",
         "针对主流程接口，覆盖相关接口的自动化平台接口测试，"
         "提升回归效率与稳定性。"),
    ]
    grid_l = Inches(0.6); grid_t = Inches(2.55)
    card_w = Inches(4.0); card_h = Inches(4.4)
    gap = Inches(0.18)
    for i, (en, cn, desc) in enumerate(items):
        left = grid_l + i * (card_w + gap)
        add_card(s, left, grid_t, card_w, card_h)
        # top accent bar
        add_rect(s, left, grid_t, card_w, Pt(2.5), fill=C.ACCENT,
                 no_line=True)
        # number
        add_text(s, left + Inches(0.45), grid_t + Inches(0.45),
                 card_w - Inches(0.9), Inches(0.4),
                 f"0{i + 1}", size=14, color=C.ACCENT, bold=True,
                 font=FONT_EN, letter_space=400)
        # english
        add_text(s, left + Inches(0.45), grid_t + Inches(0.85),
                 card_w - Inches(0.9), Inches(0.5),
                 en, size=14, color=C.TEXT_DIM, font=FONT_EN, bold=True,
                 letter_space=400)
        # cn title
        add_text(s, left + Inches(0.45), grid_t + Inches(1.35),
                 card_w - Inches(0.9), Inches(0.8),
                 cn, size=28, color=C.TEXT_PRI, bold=True)
        # divider
        add_rect(s, left + Inches(0.45), grid_t + Inches(2.2),
                 Inches(0.5), Pt(1.5), fill=C.ACCENT, no_line=True)
        # description
        add_text(s, left + Inches(0.45), grid_t + Inches(2.4),
                 card_w - Inches(0.9), Inches(1.9),
                 desc, size=13, color=C.TEXT_SEC, line_space=1.7)


def slide_12_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 12)
    add_section_eyebrow(s, "03 · WORK SUMMARY", "工作总结")

    # Card 1 - project deliveries (left, wider)
    c1_l = Inches(0.6); c1_t = Inches(2.55)
    c1_w = Inches(7.6); c1_h = Inches(4.4)
    add_card(s, c1_l, c1_t, c1_w, c1_h)
    add_rect(s, c1_l, c1_t, c1_w, Pt(2.5), fill=C.ACCENT, no_line=True)
    add_text(s, c1_l + Inches(0.45), c1_t + Inches(0.4),
             Inches(6), Inches(0.4),
             "01 · PROJECT DELIVERY", size=12, color=C.ACCENT, bold=True,
             font=FONT_EN, letter_space=400)
    add_text(s, c1_l + Inches(0.45), c1_t + Inches(0.75),
             c1_w - Inches(0.9), Inches(0.7),
             "项目按期上线交付", size=24, color=C.TEXT_PRI, bold=True)

    # sub-items
    sub_items = [
        ("UTA-dex 改造",
         [("1 月 6 日 ", C.ACCENT), ("测试完成，", C.TEXT_SEC),
          ("1 月 8 日 ", C.ACCENT), ("发布主网。", C.TEXT_SEC)]),
        ("OpenAPI 重构改造",
         [("2026-04-21 ", C.ACCENT), ("灰度发布主网；2026-05-12 请求量 ",
          C.TEXT_SEC), ("100,000 < QPS ≤ 500,000", C.ACCENT)]),
        ("出入金流程优化",
         [("V1.0 / V2.0 / V3.0 ", C.ACCENT),
          ("已全部发布主网。", C.TEXT_SEC)]),
    ]
    y = c1_t + Inches(1.75)
    for name, parts in sub_items:
        # bullet dot
        add_rect(s, c1_l + Inches(0.5), y + Inches(0.18),
                 Inches(0.12), Inches(0.12),
                 fill=C.ACCENT, no_line=True, kind=MSO_SHAPE.OVAL)
        add_text(s, c1_l + Inches(0.8), y, c1_w - Inches(1.2),
                 Inches(0.4), name, size=15, color=C.TEXT_PRI, bold=True)
        runs = [{"text": p[0], "size": 12, "color": p[1]} for p in parts]
        add_runs(s, c1_l + Inches(0.8), y + Inches(0.42),
                 c1_w - Inches(1.2), Inches(0.5),
                 [{"runs": runs}])
        y = y + Inches(0.85)

    # Card 2 - automation (right, narrower)
    c2_l = Inches(8.4); c2_t = Inches(2.55)
    c2_w = Inches(4.3); c2_h = Inches(4.4)
    add_card(s, c2_l, c2_t, c2_w, c2_h)
    add_rect(s, c2_l, c2_t, c2_w, Pt(2.5), fill=C.ACCENT_2, no_line=True)
    add_text(s, c2_l + Inches(0.45), c2_t + Inches(0.4),
             Inches(6), Inches(0.4),
             "02 · AUTOMATION", size=12, color=C.ACCENT_2, bold=True,
             font=FONT_EN, letter_space=400)
    add_text(s, c2_l + Inches(0.45), c2_t + Inches(0.75),
             c2_w - Inches(0.9), Inches(0.7),
             "自动化覆盖率达标", size=22, color=C.TEXT_PRI, bold=True)
    # big metric
    add_text(s, c2_l + Inches(0.45), c2_t + Inches(1.7),
             c2_w - Inches(0.9), Inches(1.2),
             "Local · Testnet", size=14, color=C.TEXT_DIM, font=FONT_EN,
             letter_space=200)
    add_text(s, c2_l + Inches(0.45), c2_t + Inches(1.95),
             c2_w - Inches(0.9), Inches(1.0),
             "OpenAPI", size=32, color=C.ACCENT, bold=True, font=FONT_EN)
    add_text(s, c2_l + Inches(0.45), c2_t + Inches(2.85),
             c2_w - Inches(0.9), Inches(1.8),
             "已完成 Local Testnet 环境主流程场景用例覆盖；"
             "搭建 Docker 定时任务并部署至本地试运行。",
             size=12, color=C.TEXT_SEC, line_space=1.7)


def slide_13_outlook(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 13)
    add_section_eyebrow(s, "03 · LOOKING AHEAD", "未来工作展望")

    items = [
        ("01", "提升接口测试自动化能力",
         "API AUTOMATION",
         "利用 APIFOX 或 AI 平台完善接口自动化能力，提升测试效率。"),
        ("02", "深度熟悉合约与现货业务",
         "BUSINESS DEPTH",
         "合约交易与现货交易功能继续深耕，深度理解撮合相关业务。"),
        ("03", "AI 赋能 QA 工作",
         "AI EMPOWERMENT",
         "学习使用 AI 提升测试效率：用例编写、需求分析、接口测试执行、"
         "回归验证 & 性能压测等。"),
    ]
    grid_l = Inches(0.6); grid_t = Inches(2.55)
    card_w = Inches(4.0); card_h = Inches(4.4)
    gap = Inches(0.18)
    for i, (num, cn, en, desc) in enumerate(items):
        left = grid_l + i * (card_w + gap)
        add_card(s, left, grid_t, card_w, card_h)
        add_rect(s, left, grid_t, Pt(2.5), card_h,
                 fill=C.ACCENT, no_line=True)
        add_text(s, left + Inches(0.45), grid_t + Inches(0.45),
                 Inches(2), Inches(0.55),
                 num, size=28, color=C.ACCENT, bold=True, font=FONT_EN)
        add_text(s, left + Inches(0.45), grid_t + Inches(1.1),
                 card_w - Inches(0.9), Inches(0.4),
                 en, size=11, color=C.TEXT_DIM, font=FONT_EN, bold=True,
                 letter_space=400)
        add_text(s, left + Inches(0.45), grid_t + Inches(1.5),
                 card_w - Inches(0.9), Inches(1.3),
                 cn, size=22, color=C.TEXT_PRI, bold=True,
                 line_space=1.2)
        add_rect(s, left + Inches(0.45), grid_t + Inches(2.85),
                 Inches(0.5), Pt(1.5), fill=C.ACCENT, no_line=True)
        add_text(s, left + Inches(0.45), grid_t + Inches(3.0),
                 card_w - Inches(0.9), Inches(1.3),
                 desc, size=13, color=C.TEXT_SEC, line_space=1.7)


def slide_14_suggestions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    page_chrome(s, 14)
    add_section_eyebrow(s, "03 · SUGGESTIONS", "建议")

    items = [
        ("01", "项目管理工具",
         "TOOLING",
         "工具可采购成熟的软件，如 Tapd、Jira 等；统一管理需求 / 缺陷 / "
         "迭代节奏，提升跨组协同效率。"),
        ("02", "业务告警治理",
         "ALERTING",
         "当前业务告警较为分散且频繁，建议建立告警分级与收敛机制，"
         "降低噪音、提升处理优先级判断效率。"),
    ]
    grid_l = Inches(0.6); grid_t = Inches(2.55)
    card_w = Inches(6.05); card_h = Inches(4.4)
    gap = Inches(0.2)
    for i, (num, cn, en, desc) in enumerate(items):
        left = grid_l + i * (card_w + gap)
        add_card(s, left, grid_t, card_w, card_h)
        add_rect(s, left, grid_t, Pt(2.5), card_h,
                 fill=C.ACCENT, no_line=True)
        add_text(s, left + Inches(0.6), grid_t + Inches(0.55),
                 Inches(3), Inches(0.55),
                 num, size=34, color=C.ACCENT, bold=True, font=FONT_EN)
        add_text(s, left + Inches(0.6), grid_t + Inches(1.3),
                 card_w - Inches(1.2), Inches(0.4),
                 en, size=12, color=C.TEXT_DIM, font=FONT_EN, bold=True,
                 letter_space=400)
        add_text(s, left + Inches(0.6), grid_t + Inches(1.7),
                 card_w - Inches(1.2), Inches(0.9),
                 cn, size=30, color=C.TEXT_PRI, bold=True)
        add_rect(s, left + Inches(0.6), grid_t + Inches(2.65),
                 Inches(0.6), Pt(1.5), fill=C.ACCENT, no_line=True)
        add_text(s, left + Inches(0.6), grid_t + Inches(2.85),
                 card_w - Inches(1.2), Inches(1.4),
                 desc, size=14, color=C.TEXT_SEC, line_space=1.75)


def slide_15_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # decorative
    add_glow_circle(s, 11.8, 6.8, 2.0, color=C.ACCENT, alpha=10)
    add_glow_circle(s, 1.3, 1.0, 1.5, color=C.ACCENT_2, alpha=10)
    add_grid_decor(s, Inches(10.5), Inches(0.6), cols=12, rows=4,
                   cell=Inches(0.22), color=C.ACCENT, alpha=18)
    add_top_accent(s, height_pt=3)

    # eyebrow
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             "END OF PRESENTATION",
             size=14, color=C.ACCENT, bold=True, font=FONT_EN,
             align=PP_ALIGN.CENTER, letter_space=600)
    # main THE END
    add_text(s, Inches(0.6), Inches(2.9), Inches(12.13), Inches(2.0),
             "THE  END", size=140, color=C.TEXT_PRI, bold=True,
             font=FONT_EN, align=PP_ALIGN.CENTER, letter_space=300)
    # divider
    add_rect(s, Inches(6.16), Inches(5.0), Inches(1.0), Pt(2.5),
             fill=C.ACCENT, no_line=True)
    # thanks
    add_text(s, Inches(0.6), Inches(5.2), Inches(12.13), Inches(0.8),
             "谢   谢", size=42, color=C.ACCENT_2, bold=True,
             align=PP_ALIGN.CENTER, letter_space=600)
    # footer
    add_text(s, Inches(0.6), Inches(6.7), Inches(12.13), Inches(0.4),
             "MORELL  ·  ZOOMEX QA  ·  BASE TRADING GROUP  ·  2026/05/12",
             size=11, color=C.TEXT_DIM, font=FONT_EN,
             align=PP_ALIGN.CENTER, letter_space=400)


# ============================================================
# Build
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_cover(prs)                               # 1
    slide_2_toc(prs)                                 # 2
    section_divider(prs, 1, "01",
                    "公司产品 & 文化理念",
                    "Company Products & Culture", 3)  # 3
    slide_4_culture(prs)                              # 4
    slide_5_product_web(prs)                          # 5
    slide_6_product_app(prs)                          # 6
    section_divider(prs, 2, "02",
                    "部门业务 & OKR",
                    "Department Business & OKR", 7)   # 7
    slide_8_dept(prs)                                 # 8
    slide_9_okr(prs)                                  # 9
    section_divider(prs, 3, "03",
                    "个人总结",
                    "Personal Review & Outlook", 10)  # 10
    slide_11_responsibilities(prs)                    # 11
    slide_12_summary(prs)                             # 12
    slide_13_outlook(prs)                             # 13
    slide_14_suggestions(prs)                         # 14
    slide_15_thanks(prs)                              # 15

    prs.save(OUT_PATH)
    print(f"OK -> {OUT_PATH}")


if __name__ == "__main__":
    build()
